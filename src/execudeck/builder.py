import json
import logging
from pathlib import Path
from pydantic import ValidationError
from typing import Any, Dict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

from .schema import DeckStructure, CritiqueReport
from .config import Config

logger = logging.getLogger(__name__)

# Mapping from DeckStructure layout names to python-pptx built-in blank layout index
LAYOUT_MAP = {
    "Title Slide": 0,
    "Title and Content": 1,
    "Section Header": 2,
    "Two Content": 3,
    "Comparison": 4,
    "Title Only": 5,
    "Blank": 6,
    "Content with Caption": 7,
    "Picture with Caption": 8
}

CHART_TYPE_MAP = {
    "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "COLUMN_STACKED": XL_CHART_TYPE.COLUMN_STACKED,
    "COLUMN_STACKED_100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "BAR_CLUSTERED": XL_CHART_TYPE.BAR_CLUSTERED,
    "BAR_STACKED": XL_CHART_TYPE.BAR_STACKED,
    "BAR_STACKED_100": XL_CHART_TYPE.BAR_STACKED_100,
    "LINE": XL_CHART_TYPE.LINE,
    "PIE": XL_CHART_TYPE.PIE,
    "DOUGHNUT": XL_CHART_TYPE.DOUGHNUT,
}

def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert hex string (#RRGGBB) to RGBColor object."""
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def build_presentation(deck: DeckStructure, output_path: Path, template_path: str | None = None) -> None:
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        # Clean existing slides if any
        for _ in range(len(prs.slides)):
            xml_slides = prs.slides._sldIdLst
            xml_slides.remove(xml_slides[0])

    # Store first slide's title position for stability check
    first_title_pos = None

    for i, slide_spec in enumerate(deck.slides):
        layout_idx = LAYOUT_MAP.get(slide_spec.layout, 6) # Default to Blank
        try:
            slide_layout = prs.slide_layouts[layout_idx]
        except IndexError:
            logger.warning(f"Layout index {layout_idx} not found in template. Using blank layout.")
            slide_layout = prs.slide_layouts[6]

        slide = prs.slides.add_slide(slide_layout)

        # 1. Action Title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_spec.action_title

            # Consistency check
            if first_title_pos is None:
                first_title_pos = (title_shape.left, title_shape.top)
            else:
                if (title_shape.left, title_shape.top) != first_title_pos:
                    logger.warning(f"Title position on slide {slide_spec.slide_number} differs from first slide.")
        else:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            txBox.text_frame.text = slide_spec.action_title

        # Determine if we are using Title Slide or other
        is_title_slide = slide_spec.layout == "Title Slide"

        # 2. Subtitle
        if slide_spec.subtitle:
            subtitle_shape = None
            if is_title_slide:
                 # In Title Slide, subtitle is usually index 1
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        subtitle_shape = shape
                        break

            if subtitle_shape:
                subtitle_shape.text = slide_spec.subtitle
            else:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
                txBox.text_frame.text = slide_spec.subtitle

        # 3. Body Bullets
        if slide_spec.body and slide_spec.body.bullets:
            body_shape = None
            if not is_title_slide:
                # In typical layouts like "Title and Content", body is usually index 1
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        body_shape = shape
                        break

            if not body_shape:
                body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))

            tf = body_shape.text_frame
            tf.clear()
            for bullet in slide_spec.body.bullets:
                p = tf.add_paragraph()
                p.text = bullet.text
                p.level = bullet.level
                if bullet.bold:
                    p.font.bold = True

        # 4. Chart
        if slide_spec.chart:
            chart_data = CategoryChartData()
            chart_data.categories = slide_spec.chart.categories
            for series in slide_spec.chart.series:
                chart_data.add_series(series.name, series.values)

            chart_type = CHART_TYPE_MAP.get(slide_spec.chart.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
            chart_shape = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
            chart = chart_shape.chart

            if slide_spec.chart.title:
                chart.has_title = True
                chart.chart_title.text_frame.text = slide_spec.chart.title

            # Formatting bar charts (zero enforcement)
            if chart_type in [XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED, XL_CHART_TYPE.BAR_STACKED_100]:
                if chart.value_axis:
                    chart.value_axis.minimum_scale = 0

            # Series color application and Legend removal + direct labels
            chart.has_legend = False

            for idx, series_spec in enumerate(slide_spec.chart.series):
                series_obj = chart.series[idx]
                # Direct label placement
                series_obj.has_data_labels = True

                if series_spec.color:
                    color = _hex_to_rgb(series_spec.color)
                    # For solid fills, usually we modify the fill of the entire series.
                    # Column/Bar/Pie have points, Line has line
                    if chart_type in [XL_CHART_TYPE.LINE]:
                        series_obj.format.line.color.rgb = color
                    else:
                        fill = series_obj.format.fill
                        fill.solid()
                        fill.fore_color.rgb = color

            # Annotations
            if slide_spec.chart.annotations:
                for ann in slide_spec.chart.annotations:
                    # Creating annotation text boxes over the chart area (simplified position for now)
                    txBox = slide.shapes.add_textbox(x + Inches(1), y - Inches(0.5), Inches(2), Inches(0.5))
                    txBox.text_frame.text = ann.text

        # 5. Footnotes
        if slide_spec.footnotes:
            footnote_text = "\n".join(slide_spec.footnotes)
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = footnote_text
            p.font.size = Pt(10)

        # 6. Speaker Notes
        if slide_spec.speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_spec.speaker_notes

        # 7. Page Number
        # Usually added as a slide number field, but here we add a simple text box at bottom right
        page_txBox = slide.shapes.add_textbox(Inches(9), Inches(7.0), Inches(0.5), Inches(0.5))
        page_tf = page_txBox.text_frame
        p = page_tf.paragraphs[0]
        p.text = str(slide_spec.slide_number)
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(10)

    prs.save(output_path)

def generate_markdown_report(report: CritiqueReport, output_path: Path) -> None:
    md = f"# Presentation Critique Report\n\n"
    md += f"**Overall Score:** {report.overall_score}/100\n\n"
    md += f"## Summary\n{report.summary}\n\n"

    md += "## Checklist\n"
    for section_name, section in report.checklist.model_dump().items():
        status = "✅ Pass" if section["passed"] else "❌ Fail"
        md += f"### {section['section_name']} ({status})\n"
        md += f"{section['details']}\n\n"

    md += "## Slide Details\n"
    for slide in report.slide_scores:
        md += f"### Slide {slide.slide_number} (Score: {slide.score})\n"
        if not slide.violations:
            md += "No violations.\n\n"
        else:
            for v in slide.violations:
                severity = v.severity.upper()
                md += f"- **[{severity}]** {v.description}\n"
                md += f"  *Suggested Fix:* {v.suggested_fix}\n"
        md += "\n"

    output_path.write_text(md, encoding="utf-8")

def build(json_path: str | Path, output_dir: str | Path, template_path: str | None, output_name: str, config: Config) -> Path:
    json_path = Path(json_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    try:
        data = json.loads(json_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as e:
        logger.error(f"Invalid JSON in {json_path}: {e}")
        raise ValueError("Invalid JSON file")

    # Detect type based on structure
    if "metadata" in data and "slides" in data:
        try:
            deck = DeckStructure(**data)
        except ValidationError as e:
            logger.error(f"Validation error for DeckStructure: {e}")
            raise

        output_file = output_dir / output_name
        if not output_file.suffix == ".pptx":
            output_file = output_file.with_suffix(".pptx")

        build_presentation(deck, output_file, template_path)
        logger.info(f"Built presentation: {output_file}")
        return output_file

    elif "overall_score" in data and "summary" in data:
        try:
            report = CritiqueReport(**data)
        except ValidationError as e:
            logger.error(f"Validation error for CritiqueReport: {e}")
            raise

        output_file = output_dir / output_name
        if not output_file.suffix == ".md":
            output_file = output_file.with_suffix(".md")

        generate_markdown_report(report, output_file)
        logger.info(f"Generated review report: {output_file}")
        return output_file

    else:
        raise ValueError("JSON does not match DeckStructure or CritiqueReport schema")
