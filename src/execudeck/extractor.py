"""
Deck Extractor Module.

Responsible for parsing a PowerPoint (.pptx) file to extract logical structure,
metadata, shapes, charts, and content into standard Pydantic schema representations.
Exposes the main `extract` function.
"""
import os
import logging
from pathlib import Path
from typing import List, Optional, Any, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from execudeck.schema import (
    DeckExtraction, DeckMetadata, ThemeFonts, ExtractedSlide, ExtractedShape,
    ExtractedParagraph, ExtractedChart, ExtractedSeries, ExtractedTable, ExtractedImage
)

logger = logging.getLogger(__name__)

def _is_footnote(shape: Any, slide_height: Any) -> bool:
    """
    Footnote heuristic: font size <= 12pt AND positioned in bottom 15% of the slide.
    """
    if not shape.has_text_frame:
        return False

    try:
        if shape.top is None or shape.height is None or slide_height is None:
            return False

        # Check position: bottom 15%
        bottom_edge = shape.top + shape.height
        if bottom_edge < slide_height * 0.85:
            return False

        # Check font size
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size and run.font.size > Pt(12):
                    return False
        return True
    except (AttributeError, TypeError) as e:
        logger.warning(f"Error checking footnote heuristic: {e}")
        return False

def _extract_shape_type(shape: Any) -> str:
    try:
        return shape.shape_type.name if hasattr(shape, 'shape_type') else "UNKNOWN"
    except AttributeError as e:
        logger.warning(f"Error identifying shape type: {e}")
        return "UNKNOWN"

def _extract_text_box(shape: Any, shape_type_name: str) -> Optional[ExtractedShape]:
    paragraphs = []
    for p in shape.text_frame.paragraphs:
        if not p.text.strip():
            continue

        is_bold = False
        is_italic = False
        if p.runs:
            font = p.runs[0].font
            is_bold = font.bold if font.bold is not None else False
            is_italic = font.italic if font.italic is not None else False

        paragraphs.append(ExtractedParagraph(
            text=p.text.strip(),
            level=p.level,
            bold=is_bold,
            italic=is_italic
        ))

    if paragraphs:
        return ExtractedShape(
            shape_name=shape.name,
            shape_type=shape_type_name,
            paragraphs=paragraphs
        )
    return None

def _extract_table(shape: Any) -> ExtractedTable:
    rows = []
    for row in shape.table.rows:
        rows.append([cell.text_frame.text for cell in row.cells])
    return ExtractedTable(
        shape_name=shape.name,
        rows=rows
    )

def _extract_chart(shape: Any) -> ExtractedChart:
    chart = shape.chart
    chart_type_str = str(chart.chart_type).replace("XL_CHART_TYPE.", "")
    chart_title = chart.chart_title.text_frame.text if chart.has_title else None

    categories = []
    series_data = []

    try:
        categories = [str(c.label) for c in chart.plots[0].categories]
        for s in chart.series:
            values = [v for v in s.values]
            series_data.append(ExtractedSeries(name=s.name, values=values))
    except (AttributeError, IndexError, KeyError, TypeError) as e:
        logger.warning(f"Error extracting chart data: {e}")
        pass

    return ExtractedChart(
        shape_name=shape.name,
        chart_type=chart_type_str,
        title=chart_title,
        categories=categories,
        series=series_data
    )

def _extract_image(shape: Any) -> ExtractedImage:
    return ExtractedImage(
        shape_name=shape.name,
        image_name=shape.name
    )

def _extract_slide_titles(slide: Any, placeholders: List[Any]) -> Tuple[Optional[str], Optional[str]]:
    title_text = slide.shapes.title.text if slide.shapes.title else None
    subtitle_text = None

    if len(placeholders) > 1:
        try:
             subtitle_text = placeholders[1].text if placeholders[1] != slide.shapes.title else None
        except (AttributeError, IndexError) as e:
             logger.warning(f"Error extracting subtitle: {e}")
             pass

    return title_text, subtitle_text

def _extract_speaker_notes(slide: Any) -> Optional[str]:
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
         notes_text = slide.notes_slide.notes_text_frame.text
         if notes_text.strip():
             return notes_text
    return None

def extract(pptx_path: str | Path) -> DeckExtraction:
    """Extract full content and metadata from a PowerPoint presentation."""
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"Presentation not found: {path}")

    prs = Presentation(str(path))
    slide_height = prs.slide_height

    # Basic deck metadata
    title = prs.core_properties.title if hasattr(prs.core_properties, 'title') else None
    author = prs.core_properties.author if hasattr(prs.core_properties, 'author') else None

    # Theme fonts extraction (heuristic or default as python-pptx support is limited)
    theme_fonts = ThemeFonts(major_font="Arial", minor_font="Arial")

    metadata = DeckMetadata(title=title, author=author, theme_fonts=theme_fonts)
    slides_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Blank"

        placeholders = [s for s in slide.shapes if s.is_placeholder]
        title_text, subtitle_text = _extract_slide_titles(slide, placeholders)

        body_shapes = []
        charts = []
        tables = []
        images = []
        footnotes = []

        for shape in slide.shapes:
            if _is_footnote(shape, slide_height):
                footnotes.append(shape.text)
                continue

            shape_type_name = _extract_shape_type(shape)

            if shape.has_text_frame and shape != slide.shapes.title and not (len(placeholders) > 1 and shape == placeholders[1]):
                text_box = _extract_text_box(shape, shape_type_name)
                if text_box:
                    body_shapes.append(text_box)
            elif shape.has_table:
                tables.append(_extract_table(shape))
            elif shape.has_chart:
                charts.append(_extract_chart(shape))
            elif hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append(_extract_image(shape))

        slides_data.append(ExtractedSlide(
            slide_number=idx,
            layout_name=layout_name,
            title=title_text,
            subtitle=subtitle_text,
            body_shapes=body_shapes,
            charts=charts,
            tables=tables,
            images=images,
            footnotes=footnotes,
            speaker_notes=_extract_speaker_notes(slide)
        ))

    return DeckExtraction(metadata=metadata, slides=slides_data)
