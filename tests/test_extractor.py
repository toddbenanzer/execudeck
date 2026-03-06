import pytest
import tempfile
from pathlib import Path
from pptx import Presentation
from execudeck.extractor import extract
from execudeck.schema.extraction import DeckExtraction

def test_extract_invalid_path_raises():
    with pytest.raises(FileNotFoundError):
        extract("nonexistent_path.pptx")

def test_extract_returns_deck_extraction_model(tmp_path, create_sample_pptx):
    pptx_file = tmp_path / "sample.pptx"
    create_sample_pptx(pptx_file)
    deck = extract(pptx_file)
    assert isinstance(deck, DeckExtraction)
    assert deck.metadata is not None
    assert len(deck.slides) == 3

def test_extract_title_captured(tmp_path, create_sample_pptx):
    pptx_file = tmp_path / "sample.pptx"
    create_sample_pptx(pptx_file)
    deck = extract(pptx_file)
    slide1 = deck.slides[0]
    assert slide1.title == "Hello, World!"
    assert slide1.subtitle == "Subtitle here" or slide1.subtitle is None

def test_extract_body_shapes_captured(tmp_path, create_sample_pptx):
    pptx_file = tmp_path / "sample.pptx"
    create_sample_pptx(pptx_file)
    deck = extract(pptx_file)
    slide3 = deck.slides[2]
    assert len(slide3.body_shapes) == 1
    assert slide3.body_shapes[0].paragraphs[0].text == "This is a text box"

def test_extract_chart_captured(tmp_path, create_sample_pptx):
    pptx_file = tmp_path / "sample.pptx"
    create_sample_pptx(pptx_file)
    deck = extract(pptx_file)
    slide2 = deck.slides[1]
    assert len(slide2.charts) == 1
    chart = slide2.charts[0]
    assert chart.chart_type == "BAR_CLUSTERED (57)"
    assert "A" in chart.categories
    assert len(chart.series) == 1
    assert chart.series[0].name == "Series 1"
    assert chart.series[0].values == [1.2, 2.3, 3.4]

def test_extract_footnotes_detected(tmp_path, create_sample_pptx):
    pptx_file = tmp_path / "sample.pptx"
    create_sample_pptx(pptx_file)
    deck = extract(pptx_file)
    slide3 = deck.slides[2]
    # The text should be removed from body shapes and put in footnotes
    assert len(slide3.footnotes) == 1
    assert slide3.footnotes[0] == "*1 This is a footnote"

def test_extract_speaker_notes(tmp_path, create_sample_pptx):
    pptx_file = tmp_path / "sample.pptx"
    create_sample_pptx(pptx_file)
    deck = extract(pptx_file)
    slide1 = deck.slides[0]
    assert slide1.speaker_notes == "These are speaker notes."
