import pytest
import json
from pathlib import Path
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

@pytest.fixture
def create_sample_pptx():
    def _create(path):
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_slide_layout)
        title1 = slide1.shapes.title
        subtitle1 = slide1.placeholders[1]
        title1.text = "Hello, World!"
        subtitle1.text = "Subtitle here"

        # Notes on slide 1
        notes_slide = slide1.notes_slide
        notes_slide.notes_text_frame.text = "These are speaker notes."

        # Slide 2: Data slide with a bar chart
        chart_slide_layout = prs.slide_layouts[5] # Title only
        slide2 = prs.slides.add_slide(chart_slide_layout)
        slide2.shapes.title.text = "Data Slide Title"

        chart_data = CategoryChartData()
        chart_data.categories = ['A', 'B', 'C']
        chart_data.add_series('Series 1', (1.2, 2.3, 3.4))
        slide2.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(1), Inches(2), Inches(6), Inches(4.5), chart_data
        )

        # Slide 3: Blank slide with a text box and footnote
        blank_slide_layout = prs.slide_layouts[6]
        slide3 = prs.slides.add_slide(blank_slide_layout)

        txBox = slide3.shapes.add_textbox(left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2))
        tf = txBox.text_frame
        tf.text = "This is a text box"

        # Add footnote (font size <= 12, bottom 15%)
        slide_height = prs.slide_height
        fnBox = slide3.shapes.add_textbox(left=Inches(1), top=slide_height - Inches(0.5), width=Inches(4), height=Inches(0.5))
        fn_tf = fnBox.text_frame
        p = fn_tf.paragraphs[0]
        run = p.add_run()
        run.text = "*1 This is a footnote"
        run.font.size = Pt(10)

        prs.save(path)
        return path
    return _create

@pytest.fixture
def sample_critique_json(tmp_path):
    data = {
        "slides": [
            {
                "slide_number": 1,
                "score": 90,
                "violations": [
                    {
                        "severity": "minor",
                        "description": "Topic title used instead of action title"
                    }
                ]
            }
        ],
        "checklist": {
            "story": {"passed": True, "notes": "Good"},
            "design": {"passed": True, "notes": "Good"},
            "data": {"passed": True, "notes": "Good"}
        },
        "overall_score": 90,
        "summary_notes": "Good start"
    }
    path = tmp_path / "sample_critique.json"
    with open(path, "w") as f:
        json.dump(data, f)
    return path

@pytest.fixture
def sample_deck_structure_json(tmp_path):
    data = {
        "metadata": {
            "title": "Test Deck",
            "author": "Test Author",
            "color_palette": {
                "primary": "FF0000",
                "secondary": "00FF00",
                "accent": "0000FF",
                "background": "FFFFFF",
                "text": "000000"
            }
        },
        "slides": [
            {
                "slide_number": 1,
                "layout": "title",
                "action_title": "Action Title",
                "subtitle": "Subtitle",
                "body": {"bullets": []},
                "chart": None,
                "annotations": [],
                "footnotes": [],
                "speaker_notes": "Notes"
            },
            {
                "slide_number": 2,
                "layout": "content",
                "action_title": "Chart Title",
                "subtitle": None,
                "body": {"bullets": []},
                "chart": {
                    "type": "bar",
                    "x_axis": {"label": "X", "categories": ["A", "B"]},
                    "series": [{"name": "S1", "values": [1, 2], "color_hex": "FF0000"}]
                },
                "annotations": [],
                "footnotes": [],
                "speaker_notes": ""
            },
            {
                "slide_number": 3,
                "layout": "content",
                "action_title": "Text Title",
                "subtitle": None,
                "body": {
                    "bullets": [
                        {"level": 0, "text": "Bullet 1", "bold": False}
                    ]
                },
                "chart": None,
                "annotations": [],
                "footnotes": ["*1 Footnote"],
                "speaker_notes": ""
            }
        ]
    }
    path = tmp_path / "sample_deck.json"
    with open(path, "w") as f:
        json.dump(data, f)
    return path

@pytest.fixture
def sample_content_json(tmp_path):
    data = {
        "audience": "Execs",
        "goal": "Approval",
        "key_points": ["Point 1", "Point 2"],
        "data_sources": ["Source 1"]
    }
    path = tmp_path / "sample_content.json"
    with open(path, "w") as f:
        json.dump(data, f)
    return path
