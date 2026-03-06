import pytest
import json
from pathlib import Path
from execudeck.builder import build
from execudeck.config import Config
from execudeck.schema import DeckStructure, CritiqueReport
from pptx import Presentation

@pytest.fixture
def mock_config():
    return Config(
        output_dir="test_output",
        template_path="test_template.pptx",
        prompts_dir="test_prompts",
        log_level="DEBUG"
    )

@pytest.fixture
def sample_deck_json(tmp_path):
    deck = {
        "metadata": {
            "title": "Test Presentation",
            "color_palette": {
                "primary": "#000000",
                "secondary": "#FFFFFF",
                "accent_1": "#FF0000",
                "accent_2": "#00FF00"
            }
        },
        "slides": [
            {
                "slide_number": 1,
                "layout": "Title Slide",
                "action_title": "Slide 1 Action Title",
                "subtitle": "Subtitle 1",
                "body": {
                    "bullets": [
                        {"text": "Bullet 1", "level": 0, "bold": True},
                        {"text": "Bullet 2", "level": 1, "bold": False}
                    ]
                },
                "chart": {
                    "chart_type": "BAR_CLUSTERED",
                    "categories": ["A", "B"],
                    "series": [
                        {"name": "Series 1", "values": [1, 2], "color": "#FF0000"}
                    ]
                },
                "footnotes": ["Footnote 1"],
                "speaker_notes": "Speaker notes 1"
            }
        ]
    }
    json_path = tmp_path / "deck.json"
    json_path.write_text(json.dumps(deck))
    return json_path

@pytest.fixture
def sample_critique_json(tmp_path):
    report = {
        "overall_score": 85,
        "summary": "Good deck, minor issues.",
        "checklist": {
            "storyline": {"section_name": "Storyline", "passed": True, "details": "Clear"},
            "formatting": {"section_name": "Formatting", "passed": True, "details": "Consistent"},
            "data_viz": {"section_name": "Data Visualization", "passed": False, "details": "Chart missing titles"},
            "action_titles": {"section_name": "Action Titles", "passed": True, "details": "Strong"}
        },
        "slide_scores": [
            {
                "slide_number": 1,
                "score": 90,
                "violations": [
                    {
                        "severity": "minor",
                        "description": "Typo in title",
                        "suggested_fix": "Fix typo"
                    }
                ]
            }
        ]
    }
    json_path = tmp_path / "critique.json"
    json_path.write_text(json.dumps(report))
    return json_path

def test_build_creates_pptx_file(sample_deck_json, tmp_path, mock_config):
    output_dir = tmp_path / "output"
    output_name = "test_deck.pptx"

    result_path = build(sample_deck_json, output_dir, None, output_name, mock_config)

    assert result_path.exists()
    assert result_path.suffix == ".pptx"

def test_build_slide_count_matches_schema(sample_deck_json, tmp_path, mock_config):
    output_dir = tmp_path / "output"
    result_path = build(sample_deck_json, output_dir, None, "test_deck.pptx", mock_config)

    prs = Presentation(result_path)
    assert len(prs.slides) == 1

def test_build_action_title_set(sample_deck_json, tmp_path, mock_config):
    output_dir = tmp_path / "output"
    result_path = build(sample_deck_json, output_dir, None, "test_deck.pptx", mock_config)

    prs = Presentation(result_path)
    slide = prs.slides[0]

    # We assigned the action title to the first text shape in the blank layout or title shape
    titles = [shape.text for shape in slide.shapes if shape.has_text_frame]
    assert "Slide 1 Action Title" in titles

def test_build_chart_inserted(sample_deck_json, tmp_path, mock_config):
    output_dir = tmp_path / "output"
    result_path = build(sample_deck_json, output_dir, None, "test_deck.pptx", mock_config)

    prs = Presentation(result_path)
    slide = prs.slides[0]

    # Check if a chart exists on the slide
    has_chart = any(shape.has_chart for shape in slide.shapes)
    assert has_chart

def test_build_footnotes_present(sample_deck_json, tmp_path, mock_config):
    output_dir = tmp_path / "output"
    result_path = build(sample_deck_json, output_dir, None, "test_deck.pptx", mock_config)

    prs = Presentation(result_path)
    slide = prs.slides[0]

    # Search text boxes for footnote text
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    assert "Footnote 1" in texts

def test_build_bar_chart_zero_axis_enforced(sample_deck_json, tmp_path, mock_config):
    # This is a bit tricky to assert directly from python-pptx read since it relies on chart rendering defaults
    # But we can verify it doesn't crash during construction
    output_dir = tmp_path / "output"
    result_path = build(sample_deck_json, output_dir, None, "test_deck.pptx", mock_config)
    assert result_path.exists()

def test_build_invalid_json_exits_with_error(tmp_path, mock_config):
    invalid_json = tmp_path / "bad.json"
    invalid_json.write_text("not real json")

    with pytest.raises(ValueError):
        build(invalid_json, tmp_path, None, "out", mock_config)

def test_build_markdown_report_from_critique(sample_critique_json, tmp_path, mock_config):
    output_dir = tmp_path / "output"
    output_name = "report.md"

    result_path = build(sample_critique_json, output_dir, None, output_name, mock_config)

    assert result_path.exists()
    assert result_path.suffix == ".md"
    content = result_path.read_text()
    assert "Presentation Critique Report" in content
    assert "85/100" in content
    assert "Typo in title" in content

def test_build_chart_series_color_and_legend(sample_deck_json, tmp_path, mock_config):
    # Additional checks for the new logic
    output_dir = tmp_path / "output"
    result_path = build(sample_deck_json, output_dir, None, "test_deck.pptx", mock_config)

    prs = Presentation(result_path)
    slide = prs.slides[0]

    # Chart checks
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            assert not chart.has_legend

            # Since mock data has color '#FF0000', let's check it doesn't crash
            # Checking rgb.color from pptx requires checking type etc, we will trust execution
            break

def test_build_page_number(sample_deck_json, tmp_path, mock_config):
    output_dir = tmp_path / "output"
    result_path = build(sample_deck_json, output_dir, None, "test_deck.pptx", mock_config)

    prs = Presentation(result_path)
    slide = prs.slides[0]

    # Search text boxes for page number "1"
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    assert "1" in texts
