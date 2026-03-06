import json
import subprocess
from pathlib import Path
from pptx import Presentation

def create_sample_files():
    # 1. pptx
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Sample Presentation"
    slide.placeholders[1].text = "Subtitle"
    prs.save("e2e_tests/sample.pptx")

    # 2. content.json
    Path("e2e_tests/content.json").write_text(json.dumps({
        "objective": "Test Generate",
        "audience": "Developers",
        "content_slides": [
            {
                "key_message": "Automation saves time",
                "supporting_points": ["Point A", "Point B"]
            }
        ]
    }))

    # 3. deck_structure.json
    Path("e2e_tests/deck_structure.json").write_text(json.dumps({
        "metadata": {
            "title": "Built Deck",
            "color_palette": {"primary": "#000000", "secondary": "#FFFFFF", "accent_1": "#FF0000", "accent_2": "#00FF00"}
        },
        "slides": [
            {
                "slide_number": 1,
                "layout": "Title Slide",
                "action_title": "Automated Build Slide",
                "body": {"bullets": [{"text": "It works!", "level": 0, "bold": True}]}
            }
        ]
    }))

    # 4. critique.json
    Path("e2e_tests/critique.json").write_text(json.dumps({
        "overall_score": 90,
        "summary": "Great mock deck.",
        "checklist": {
            "storyline": {"section_name": "Storyline", "passed": True, "details": "Clear"},
            "formatting": {"section_name": "Formatting", "passed": True, "details": "Consistent"},
            "data_viz": {"section_name": "Data Viz", "passed": True, "details": "Good"},
            "action_titles": {"section_name": "Action Titles", "passed": True, "details": "Strong"}
        },
        "slide_scores": [
            {
                "slide_number": 1,
                "score": 90,
                "violations": [{"severity": "minor", "description": "Fix font", "suggested_fix": "Change font"}]
            }
        ]
    }))

def run_command(cmd):
    print(f"Running: {cmd}")
    result = subprocess.run(cmd, shell=True, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"Error executing {cmd}")
        print(result.stderr)
        exit(1)
    print(result.stdout)

def main():
    create_sample_files()

    # Test Review
    run_command("execudeck review e2e_tests/sample.pptx --output-dir e2e_tests/output/review")

    # Test Generate
    run_command("execudeck generate e2e_tests/content.json --output-dir e2e_tests/output/generate")

    # Test Edit
    run_command("execudeck edit e2e_tests/sample.pptx e2e_tests/critique.json --output-dir e2e_tests/output/edit")

    # Test Build (PPTX)
    run_command("execudeck build e2e_tests/deck_structure.json --output-dir e2e_tests/output/build --output-name out_deck.pptx")

    # Test Build (Critique Report)
    run_command("execudeck build e2e_tests/critique.json --output-dir e2e_tests/output/build --output-name out_report.md")

    print("All End-to-End tests passed.")

if __name__ == "__main__":
    main()
